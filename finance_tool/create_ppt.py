"""
生成财务自动化工具展示PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


def add_background(slide, color=RGBColor(0xF0, 0xF2, 0xF6)):
    """添加背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, left, top, width, height, font_size=28, bold=True, color=RGBColor(0x1F, 0x77, 0xB4)):
    """添加标题形状"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox


def add_content_shape(slide, text, left, top, width, height, font_size=14, color=RGBColor(0x33, 0x33, 0x33)):
    """添加内容形状"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    
    return txBox


def add_icon_card(slide, icon, title, desc, left, top, width=Inches(2.2), height=Inches(1.5)):
    """添加图标卡片"""
    # 背景框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    
    # 图标
    icon_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), Inches(0.5), Inches(0.5))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(24)
    
    # 标题
    title_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.6), width - Inches(0.2), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
    
    # 描述
    desc_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.95), width - Inches(0.2), Inches(0.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def create_presentation():
    """创建PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ==================== 第1页：封面 ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_background(slide1, RGBColor(0x1F, 0x77, 0xB4))
    
    # 主标题
    add_title_shape(
        slide1, "财务自动化工具",
        Inches(1), Inches(2), Inches(11), Inches(1.5),
        font_size=44, color=RGBColor(0xFF, 0xFF, 0xFF)
    )
    
    # 副标题
    add_content_shape(
        slide1, "Python Streamlit + openpyxl + Excel VBA",
        Inches(1), Inches(3.5), Inches(11), Inches(0.8),
        font_size=20, color=RGBColor(0xCC, 0xDD, 0xEE)
    )
    
    # 核心功能标签
    features = ["CSV导入", "统计分析", "数据可视化", "异常检测", "报表导出"]
    for i, feat in enumerate(features):
        shape = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1 + i * 2.2), Inches(4.8), Inches(2), Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
    
    # ==================== 第2页：数据导入 ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide2)
    
    add_title_shape(
        slide2, "数据导入与预览",
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
        font_size=32
    )
    
    # 功能卡片
    cards = [
        ("📁", "多格式支持", "CSV文件上传\nUTF-8/GBK/GB2312编码\n逗号/分号/制表符分隔"),
        ("✅", "智能验证", "自动检测数据类型\n缺失值统计\n行列信息展示"),
        ("📊", "数据预览", "前100行快速预览\n列类型识别\n基本统计信息"),
        ("⚙️", "列映射配置", "日期列/金额列指定\n分类列映射\n配置自动保存"),
    ]
    
    for i, (icon, title, desc) in enumerate(cards):
        add_icon_card(slide2, icon, title, desc, Inches(0.5 + i * 3.1), Inches(1.5))
    
    # 底部说明
    add_content_shape(
        slide2, "支持示例数据：上市公司年报（12家公司、6个行业、16项财务指标）",
        Inches(0.5), Inches(5.5), Inches(12), Inches(0.6),
        font_size=12, color=RGBColor(0x66, 0x66, 0x66)
    )
    
    # ==================== 第3页：统计分析 ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide3)
    
    add_title_shape(
        slide3, "统计分析与可视化",
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
        font_size=32
    )
    
    # 左侧：统计指标
    add_title_shape(
        slide3, "核心统计指标",
        Inches(0.5), Inches(1.3), Inches(5), Inches(0.5),
        font_size=18, color=RGBColor(0x33, 0x33, 0x33)
    )
    
    stats = ["均值、中位数、众数", "标准差、方差", "最大值、最小值", "四分位数", "偏度、峰度", "求和、计数"]
    for i, stat in enumerate(stats):
        y = Inches(2 + i * 0.55)
        shape = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y, Inches(5), Inches(0.45)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        p = tf.paragraphs[0]
        p.text = f"  • {stat}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # 右侧：图表类型
    add_title_shape(
        slide3, "可视化图表",
        Inches(6.5), Inches(1.3), Inches(6), Inches(0.5),
        font_size=18, color=RGBColor(0x33, 0x33, 0x33)
    )
    
    charts = [
        ("柱状图", "对比分析"),
        ("折线图", "趋势展示"),
        ("饼图", "占比分析"),
        ("箱线图", "分布识别"),
        ("散点图", "相关分析"),
    ]
    
    for i, (chart, desc) in enumerate(charts):
        y = Inches(2 + i * 0.8)
        
        # 图表名
        box = slide3.shapes.add_textbox(Inches(6.5), y, Inches(2.5), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = chart
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        
        # 描述
        box2 = slide3.shapes.add_textbox(Inches(9), y, Inches(3), Inches(0.4))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    # ==================== 第4页：高级分析 ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide4)
    
    add_title_shape(
        slide4, "智能分析引擎",
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
        font_size=32
    )
    
    # 三大分析模块
    modules = [
        ("异常检测", [
            "基于标准差的异常识别",
            "可调节检测阈值(1-4倍)",
            "异常值可视化标记",
            "偏离程度量化分析"
        ]),
        ("分组汇总", [
            "按部门/科目/类别分组",
            "金额合计与占比计算",
            "多维度交叉分析",
            "自动排序与筛选"
        ]),
        ("趋势分析", [
            "时间序列数据处理",
            "环比/同比增长率",
            "相关性热力图",
            "预测趋势线"
        ]),
    ]
    
    for i, (title, items) in enumerate(modules):
        x = Inches(0.5 + i * 4.2)
        
        # 模块标题框
        shape = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.5), Inches(3.8), Inches(4.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        shape.line.width = Pt(2)
        
        # 标题
        title_box = slide4.shapes.add_textbox(x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        p.alignment = PP_ALIGN.CENTER
        
        # 分隔线
        line = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.5), Inches(2.3), Inches(2.8), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        line.line.fill.background()
        
        # 功能列表
        for j, item in enumerate(items):
            item_box = slide4.shapes.add_textbox(x + Inches(0.3), Inches(2.6 + j * 0.6), Inches(3.2), Inches(0.5))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # ==================== 第5页：报表导出 ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide5)
    
    add_title_shape(
        slide5, "一键报表导出",
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
        font_size=32
    )
    
    # 左侧：Excel导出
    add_title_shape(
        slide5, "Excel报表",
        Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.5),
        font_size=20, color=RGBColor(0x1F, 0x77, 0xB4)
    )
    
    excel_features = [
        "原始数据表：完整数据备份",
        "统计汇总表：核心指标计算",
        "分类汇总表：按类别分组统计",
        "数据透视表：多维度交叉分析",
        "内嵌图表：柱状图、饼图、折线图"
    ]
    
    for i, feat in enumerate(excel_features):
        y = Inches(2 + i * 0.65)
        shape = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y, Inches(5.5), Inches(0.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
        shape.line.fill.background()
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"  ✓ {feat}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # 右侧：VBA导出
    add_title_shape(
        slide5, "VBA模板报表",
        Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
        font_size=20, color=RGBColor(0x1F, 0x77, 0xB4)
    )
    
    vba_features = [
        "自动生成财务综合报表",
        "部门汇总与占比分析",
        "内嵌柱状图与饼图",
        "支持导出为PDF格式",
        "一键刷新数据功能"
    ]
    
    for i, feat in enumerate(vba_features):
        y = Inches(2 + i * 0.65)
        shape = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7), y, Inches(5.5), Inches(0.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
        shape.line.fill.background()
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"  ✓ {feat}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # ==================== 第6页：技术架构 ====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide6, RGBColor(0x1F, 0x77, 0xB4))
    
    add_title_shape(
        slide6, "技术架构",
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
        font_size=32, color=RGBColor(0xFF, 0xFF, 0xFF)
    )
    
    # 技术栈
    tech_stack = [
        ("Python 3.8+", "核心运行环境"),
        ("Streamlit", "Web界面框架"),
        ("Pandas", "数据处理引擎"),
        ("openpyxl", "Excel操作库"),
        ("Plotly", "可视化引擎"),
        ("Excel VBA", "报表自动化"),
    ]
    
    for i, (tech, desc) in enumerate(tech_stack):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.5 + row * 2.5)
        
        # 技术卡片
        shape = slide6.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(3.8), Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.fill.background()
        
        # 技术名
        tech_box = slide6.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.6))
        tf = tech_box.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide6.shapes.add_textbox(x + Inches(0.2), y + Inches(1.1), Inches(3.4), Inches(0.6))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    output_path = os.path.join(os.path.dirname(__file__), "财务自动化工具介绍.pptx")
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
