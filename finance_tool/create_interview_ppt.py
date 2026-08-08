"""
生成面试演示文稿 - 财务自动化工具
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os


def create_interview_ppt():
    """创建面试演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    BLUE = RGBColor(0x1F, 0x77, 0xB4)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
    
    # ==================== 第1页：项目概述 ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    
    # 标题
    box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "财务自动化分析工具"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 解决的问题
    box2 = slide1.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    lines = [
        "痛点：财务人员手动处理Excel效率低、易出错、重复劳动多",
        "目标：零代码可视化平台，让非技术人员自助完成财务分析"
    ]
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
    
    # 技术栈标签
    techs = ["Python", "Streamlit", "Pandas", "openpyxl", "Plotly", "VBA"]
    for i, tech in enumerate(techs):
        shape = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5 + i * 1.8), Inches(5.5), Inches(1.6), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE
    
    # ==================== 第2页：架构设计 ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide2.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    
    box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "系统架构"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # 架构层次
    layers = [
        ("展示层", "Streamlit Web UI", "交互式界面、图表展示、表单控件"),
        ("业务层", "Pandas + NumPy", "数据清洗、统计计算、异常检测"),
        ("导出层", "openpyxl + VBA", "Excel生成、样式渲染、宏脚本注入"),
    ]
    
    for i, (name, tech, desc) in enumerate(layers):
        y = Inches(1.5 + i * 1.8)
        
        # 层名称
        shape = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y, Inches(2.5), Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 技术
        box2 = slide2.shapes.add_textbox(Inches(3.5), y, Inches(3), Inches(1.3))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = tech
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = DARK
        p2.space_before = Pt(15)
        
        # 描述
        box3 = slide2.shapes.add_textbox(Inches(6.5), y, Inches(6), Inches(1.3))
        tf3 = box3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY
        p3.space_before = Pt(18)
    
    # ==================== 第3页：核心功能 ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide3.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    
    box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心功能"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    features = [
        ("数据导入", "CSV多编码支持\n智能列识别\n数据验证"),
        ("统计分析", "均值/标准差/极值\n多列并行计算\n分组汇总"),
        ("可视化", "5种图表类型\nPlotly交互式\n一键切换"),
        ("异常检测", "标准差阈值\n可视化标记\n偏离度量化"),
        ("报表导出", "Excel多Sheet\nVBA模板注入\nPDF导出"),
    ]
    
    for i, (title, desc) in enumerate(features):
        x = Inches(0.3 + i * 2.6)
        
        shape = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.5), Inches(2.4), Inches(4.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BLUE
        shape.line.width = Pt(2)
        
        # 标题
        box2 = slide3.shapes.add_textbox(x + Inches(0.1), Inches(1.8), Inches(2.2), Inches(0.5))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = BLUE
        p2.alignment = PP_ALIGN.CENTER
        
        # 分隔线
        line = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.3), Inches(2.4), Inches(1.8), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE
        line.line.fill.background()
        
        # 描述
        box3 = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.7), Inches(2), Inches(3))
        tf3 = box3.text_frame
        tf3.word_wrap = True
        for j, line_text in enumerate(desc.split('\n')):
            p = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
            p.text = f"• {line_text}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.space_after = Pt(8)
    
    # ==================== 第4页：技术亮点 ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide4.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    
    box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "技术亮点"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    highlights = [
        ("模块化设计", "数据处理、Excel操作、配置管理三层解耦，职责单一，便于维护扩展"),
        ("配置持久化", "JSON配置文件自动保存用户偏好，列映射、导出设置一次配置反复使用"),
        ("VBA深度集成", "Python生成Excel后注入VBA宏脚本，实现跨语言自动化报表流水线"),
        ("鲁棒性处理", "多编码自动探测、异常值智能识别、数据验证机制，保障数据质量"),
    ]
    
    for i, (title, desc) in enumerate(highlights):
        y = Inches(1.3 + i * 1.4)
        
        shape = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y, Inches(12), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.fill.background()
        
        # 序号
        num_box = slide4.shapes.add_textbox(Inches(0.8), y + Inches(0.15), Inches(0.6), Inches(0.6))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(i + 1)
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = BLUE
        
        # 标题
        title_box = slide4.shapes.add_textbox(Inches(1.6), y + Inches(0.1), Inches(3), Inches(0.5))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK
        
        # 描述
        desc_box = slide4.shapes.add_textbox(Inches(1.6), y + Inches(0.6), Inches(10.5), Inches(0.5))
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = GRAY
    
    # ==================== 第5页：项目成果 ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide5.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    
    box = slide5.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "项目成果"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 成果卡片
    results = [
        ("效率提升", "报表生成时间\n从2小时→5分钟"),
        ("零门槛", "非技术人员\n可独立操作"),
        ("可扩展", "模块化架构\n易于功能迭代"),
    ]
    
    for i, (title, desc) in enumerate(results):
        x = Inches(1 + i * 4)
        
        shape = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(2.5), Inches(3.5), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.fill.background()
        
        # 标题
        box2 = slide5.shapes.add_textbox(x + Inches(0.2), Inches(2.8), Inches(3.1), Inches(0.6))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = BLUE
        p2.alignment = PP_ALIGN.CENTER
        
        # 描述
        box3 = slide5.shapes.add_textbox(x + Inches(0.2), Inches(3.5), Inches(3.1), Inches(1.2))
        tf3 = box3.text_frame
        tf3.word_wrap = True
        for j, line in enumerate(desc.split('\n')):
            p = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(5)
    
    # 总结
    box4 = slide5.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11), Inches(1))
    tf4 = box4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "技术栈：Python + Streamlit + Pandas + openpyxl + Plotly + Excel VBA"
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p4.alignment = PP_ALIGN.CENTER
    
    # 保存
    output_path = os.path.join(os.path.dirname(__file__), "面试演示-财务自动化工具.pptx")
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    create_interview_ppt()
